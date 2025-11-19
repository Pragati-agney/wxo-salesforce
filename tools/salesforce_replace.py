"""
Salesforce File Download Tool for watsonx Orchestrate

This tool retrieves files from Salesforce, modifies PowerPoint templates with user data,
and returns them as bytes for download.This is test.
"""

import logging
import requests
from requests.exceptions import RequestException, Timeout, HTTPError
from pydantic import BaseModel, Field
from typing import Literal, Optional
from datetime import datetime
from io import BytesIO
from pptx import Presentation

# Tool decorator - will be provided by watsonx orchestrate at runtime
try:
    from ibm_watsonx_orchestrate.agent_builder.tools import tool
    from ibm_watsonx_orchestrate.agent_builder.connections import ConnectionType
    from ibm_watsonx_orchestrate.run import connections
except ImportError:
    # Fallback for local testing
    def tool(*args, **kwargs):
        def decorator(func):
            return func
        return decorator
    
    class ConnectionType:
        OAUTH2_AUTH_CODE = "oauth2_auth_code"

# Configure logging
logger = logging.getLogger(__name__)

# Salesforce connection configuration
SALESFORCE_APP_ID = 'salesforce-wxo'
SALESFORCE_API_VERSION = 'v58.0'
REQUEST_TIMEOUT = 60


class SalesforceFileInput(BaseModel):
    """Input model for Salesforce file download with template modification"""
    file_id: str = Field(
        default="069NS00000VGfOjYAL",
        description="The Salesforce file ID. Supports ContentDocument (069), ContentVersion (068), or Attachment (00P) IDs"
    )
    company_name: str = Field(
        default="Acme Corporation",
        description="Company name to replace <Company> placeholder in the template"
    )
    tier: str = Field(
        default="Gold",
        description="Tier level to replace <Tier> placeholder (e.g., Silver, Gold, Platinum)"
    )


def modify_pptx_template(file_content: bytes, company_name: str, tier: str) -> bytes:
    """
    Modify PowerPoint template by replacing placeholders and updating valid through date.
    
    Args:
        file_content: Original PowerPoint file as bytes
        company_name: Company name to insert
        tier: Tier level to insert
        
    Returns:
        Modified PowerPoint file as bytes
    """
    try:
        # Load the PowerPoint from bytes
        prs = Presentation(BytesIO(file_content))
        
        # Calculate the valid through date (31 December of current year)
        current_year = datetime.now().year
        valid_through_date = f"31 December {current_year}"
        
        logger.info(f"Modifying template: Company={company_name}, Tier={tier}, Valid through={valid_through_date}")
        
        # Iterate through all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Process each paragraph and run
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Replace <Company> placeholder
                            if "<Company>" in run.text:
                                run.text = run.text.replace("<Company>", company_name)
                                logger.info(f"Replaced <Company> with {company_name}")
                            
                            # Replace <Tier> placeholder
                            if "<Tier>" in run.text:
                                run.text = run.text.replace("<Tier>", tier)
                                logger.info(f"Replaced <Tier> with {tier}")
                            
                            # Update the Valid through date (find and replace old year)
                            if "31 December" in run.text and "Valid through" in shape.text:
                                # Replace any year pattern (20XX or 19XX)
                                import re
                                run.text = re.sub(r'31 December \d{4}', valid_through_date, run.text)
                                logger.info(f"Updated Valid through date to {valid_through_date}")
        
        # Save modified presentation to bytes
        output_stream = BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        modified_content = output_stream.read()
        
        logger.info(f"Successfully modified PowerPoint template ({len(modified_content)} bytes)")
        return modified_content
        
    except Exception as e:
        logger.error(f"Error modifying PowerPoint template: {str(e)}")
        raise


@tool(
    expected_credentials=[
        {"app_id": SALESFORCE_APP_ID, "type": ConnectionType.OAUTH2_AUTH_CODE}
    ]
)
def salesforce_download_file(input_data: SalesforceFileInput) -> bytes:
    """
    Download a file from Salesforce, modify it with user data, and return it as bytes.
    
    This tool retrieves PowerPoint certificate templates from Salesforce, replaces placeholders
    with user-provided data, and calculates the validity date automatically.
    
    Replacements made:
    - <Company> → User-provided company name
    - <Tier> → User-provided tier level (Silver, Gold, Platinum, etc.)
    - Valid through date → Automatically set to 31 December of the current year
    
    Use this tool when a user requests:
    - To generate a partner certificate with their company details
    - To download a customized certificate from Salesforce
    - To create a certificate with updated information
    
    The tool returns the modified file as bytes that can be downloaded by the user.
    
    Args:
        input_data: Contains file_id, company_name, and tier
        
    Returns:
        bytes: The modified file content as bytes for download
    """
    
    try:
        # Extract inputs
        file_id = input_data.file_id.strip()
        company_name = input_data.company_name.strip()
        tier = input_data.tier.strip()
        
        # Validate file ID
        if not file_id:
            raise ValueError("file_id cannot be empty")
        
        # Get OAuth2 credentials from the connection
        creds = connections.oauth2_auth_code(SALESFORCE_APP_ID)
        base_url = creds.url.rstrip('/')
        
        # Prepare authorization headers
        headers = {
            "Authorization": f"Bearer {creds.access_token}",
            "Accept": "*/*"
        }
        
        # Handle ContentDocument ID (069) - need to get ContentVersion ID first
        if file_id.startswith('069'):
            logger.info(f"ContentDocument ID detected: {file_id}")
            
            # Query for the latest ContentVersion
            query = (
                f"SELECT Id FROM ContentVersion "
                f"WHERE ContentDocumentId = '{file_id}' AND IsLatest = true"
            )
            query_url = f"{base_url}/services/data/{SALESFORCE_API_VERSION}/query"
            
            query_response = requests.get(
                query_url, 
                headers=headers, 
                params={"q": query}, 
                timeout=REQUEST_TIMEOUT
            )
            query_response.raise_for_status()
            
            query_data = query_response.json()
            
            if not query_data.get('records'):
                raise ValueError(
                    f"No file found with ContentDocument ID: {file_id}"
                )
            
            # Get the ContentVersion ID
            content_version_id = query_data['records'][0]['Id']
            file_id = content_version_id
            logger.info(f"Found ContentVersion ID: {file_id}")
            
            # Set download URL for ContentVersion
            download_url = (
                f"{base_url}/services/data/{SALESFORCE_API_VERSION}/"
                f"sobjects/ContentVersion/{file_id}/VersionData"
            )
        
        # Handle ContentVersion ID (068)
        elif file_id.startswith('068'):
            logger.info(f"ContentVersion ID detected: {file_id}")
            download_url = (
                f"{base_url}/services/data/{SALESFORCE_API_VERSION}/"
                f"sobjects/ContentVersion/{file_id}/VersionData"
            )
        
        # Handle Attachment ID (00P)
        elif file_id.startswith('00P'):
            logger.info(f"Attachment ID detected: {file_id}")
            download_url = (
                f"{base_url}/services/data/{SALESFORCE_API_VERSION}/"
                f"sobjects/Attachment/{file_id}/Body"
            )
        
        else:
            raise ValueError(
                f"Unknown Salesforce ID format: {file_id}. "
                f"Expected ContentDocument (069), ContentVersion (068), "
                f"or Attachment (00P)"
            )
        
        # Download the file with streaming for large files
        logger.info(f"Downloading file from: {download_url}")
        
        response = requests.get(
            download_url,
            headers=headers,
            timeout=REQUEST_TIMEOUT,
            stream=True
        )
        
        # Raise exception for HTTP errors
        response.raise_for_status()
        
        # Get file content as bytes
        file_content = response.content
        file_size = len(file_content)
        
        logger.info(f"Successfully retrieved file ({file_size} bytes)")
        
        # Modify the PowerPoint template
        modified_content = modify_pptx_template(file_content, company_name, tier)
        
        # Return modified bytes for download
        return modified_content
        
    except Exception as e:
        # Return error message as bytes (text file)
        error_msg = f"Error processing file from Salesforce: {str(e)}"
        logger.error(error_msg)
        return error_msg.encode('utf-8')


# Example usage for testing
if __name__ == "__main__":
    # Test the function locally
    test_input = SalesforceFileInput(
        file_id="069NS00000VGfOjYAL",
        company_name="Tech Innovations Inc.",
        tier="Platinum"
    )
    
    # Note: This will fail without proper Salesforce credentials
    # Set environment variables for local testing:
    # export SALESFORCE_WXO_URL="https://your-instance.salesforce.com"
    # export SALESFORCE_WXO_ACCESS_TOKEN="your_access_token"
    
    print("Testing Salesforce file download with template modification...")
    print(f"File ID: {test_input.file_id}")
    print(f"Company: {test_input.company_name}")
    print(f"Tier: {test_input.tier}")
    print(f"Valid through: 31 December {datetime.now().year}")