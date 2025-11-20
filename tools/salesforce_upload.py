"""
Salesforce File Download and Upload Tool for watsonx Orchestrate

This tool retrieves files from Salesforce, modifies PowerPoint templates with user data,
uploads the modified file back to Salesforce, and returns the file as bytes.
"""

import logging
import requests
from requests.exceptions import RequestException, Timeout, HTTPError
from pydantic import BaseModel, Field
from typing import Literal, Optional
from datetime import datetime
from io import BytesIO
from pptx import Presentation
import mimetypes

# Tool decorator - will be provided by watsonx orchestrate at runtime
try:
    from ibm_watsonx_orchestrate.agent_builder.tools import tool
    from ibm_watsonx_orchestrate.agent_builder.connections import ConnectionType
    from ibm_watsonx_orchestrate.run import connections
except ImportError:
    # Fallback for local testing
    def tool(*args, **kwargs):
        def decorator(func):
            return func
        return decorator
    
    class ConnectionType:
        OAUTH2_AUTH_CODE = "oauth2_auth_code"

# Configure logging
logger = logging.getLogger(__name__)

# Salesforce connection configuration
SALESFORCE_APP_ID = 'salesforce-wxo'
SALESFORCE_API_VERSION = 'v58.0'
REQUEST_TIMEOUT = 60


class SalesforceFileInput(BaseModel):
    """Input model for Salesforce file download with template modification"""
    file_id: str = Field(
        default="069NS00000VGfOjYAL",
        description="The Salesforce file ID to download. Supports ContentDocument (069), ContentVersion (068), or Attachment (00P) IDs"
    )
    company_name: str = Field(
        default="Acme Corporation",
        description="Company name to replace <Company> placeholder in the template"
    )
    tier: str = Field(
        default="Gold",
        description="Tier level to replace <Tier> placeholder (e.g., Silver, Gold, Platinum)"
    )
    upload_back_to_salesforce: bool = Field(
        default=True,
        description="If True, upload the modified file back to Salesforce as a new version"
    )
    title: str = Field(
        default="Partner_Plus_Certificate",
        description="Title for the uploaded file in Salesforce"
    )


def modify_pptx_template(file_content: bytes, company_name: str, tier: str) -> bytes:
    """
    Modify PowerPoint template by replacing placeholders and updating valid through date.
    
    Args:
        file_content: Original PowerPoint file as bytes
        company_name: Company name to insert
        tier: Tier level to insert
        
    Returns:
        Modified PowerPoint file as bytes
    """
    try:
        # Load the PowerPoint from bytes
        prs = Presentation(BytesIO(file_content))
        
        # Calculate the valid through date (31 December of current year)
        current_year = datetime.now().year
        valid_through_date = f"31 December {current_year}"
        
        logger.info(f"Modifying template: Company={company_name}, Tier={tier}, Valid through={valid_through_date}")
        
        # Iterate through all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Process each paragraph and run
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Replace <Company> placeholder
                            if "<Company>" in run.text:
                                run.text = run.text.replace("<Company>", company_name)
                                logger.info(f"Replaced <Company> with {company_name}")
                            
                            # Replace <Tier> placeholder
                            if "<Tier>" in run.text:
                                run.text = run.text.replace("<Tier>", tier)
                                logger.info(f"Replaced <Tier> with {tier}")
                            
                            # Update the Valid through date (find and replace old year)
                            if "31 December" in run.text and "Valid through" in shape.text:
                                # Replace any year pattern (20XX or 19XX)
                                import re
                                run.text = re.sub(r'31 December \d{4}', valid_through_date, run.text)
                                logger.info(f"Updated Valid through date to {valid_through_date}")
        
        # Save modified presentation to bytes
        output_stream = BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        modified_content = output_stream.read()
        
        logger.info(f"Successfully modified PowerPoint template ({len(modified_content)} bytes)")
        return modified_content
        
    except Exception as e:
        logger.error(f"Error modifying PowerPoint template: {str(e)}")
        raise


def upload_file_to_salesforce(
    base_url: str,
    access_token: str,
    file_content: bytes,
    title: str,
    original_content_document_id: Optional[str] = None
) -> dict:
    """
    Upload a file to Salesforce as a new ContentVersion.
    
    Args:
        base_url: Salesforce instance URL
        access_token: OAuth2 access token
        file_content: File content as bytes
        title: Title for the file
        original_content_document_id: Optional ContentDocument ID to create a new version of existing file
        
    Returns:
        Dictionary with upload details including ContentVersion ID and ContentDocument ID
    """
    try:
        # Prepare the file upload
        upload_url = f"{base_url}/services/data/{SALESFORCE_API_VERSION}/sobjects/ContentVersion"
        
        headers = {
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json"
        }
        
        # Encode file content to base64
        import base64
        encoded_file = base64.b64encode(file_content).decode('utf-8')
        
        # Prepare the payload
        payload = {
            "Title": title,
            "PathOnClient": f"{title}.pptx",
            "VersionData": encoded_file,
            "IsMajorVersion": True
        }
        
        # If we have an original ContentDocument ID, link to it (create new version)
        if original_content_document_id:
            payload["ContentDocumentId"] = original_content_document_id
            logger.info(f"Creating new version for ContentDocument: {original_content_document_id}")
        else:
            logger.info("Creating new ContentDocument")
        
        # Upload the file
        logger.info(f"Uploading file to Salesforce: {title}")
        response = requests.post(
            upload_url,
            headers=headers,
            json=payload,
            timeout=REQUEST_TIMEOUT
        )
        
        response.raise_for_status()
        upload_result = response.json()
        
        if upload_result.get('success'):
            content_version_id = upload_result['id']
            logger.info(f"Successfully uploaded file. ContentVersion ID: {content_version_id}")
            
            # Query to get the ContentDocument ID
            query = f"SELECT ContentDocumentId FROM ContentVersion WHERE Id = '{content_version_id}'"
            query_url = f"{base_url}/services/data/{SALESFORCE_API_VERSION}/query"
            query_response = requests.get(
                query_url,
                headers={"Authorization": f"Bearer {access_token}"},
                params={"q": query},
                timeout=REQUEST_TIMEOUT
            )
            query_response.raise_for_status()
            query_data = query_response.json()
            
            content_document_id = None
            if query_data.get('records'):
                content_document_id = query_data['records'][0]['ContentDocumentId']
            
            return {
                "success": True,
                "content_version_id": content_version_id,
                "content_document_id": content_document_id,
                "title": title,
                "message": f"File uploaded successfully to Salesforce"
            }
        else:
            raise ValueError(f"Upload failed: {upload_result}")
            
    except Exception as e:
        logger.error(f"Error uploading file to Salesforce: {str(e)}")
        raise


@tool(
    expected_credentials=[
        {"app_id": SALESFORCE_APP_ID, "type": ConnectionType.OAUTH2_AUTH_CODE}
    ]
)
def salesforce_download_file(input_data: SalesforceFileInput) -> bytes:
    """
    Download a file from Salesforce, modify it with user data, optionally upload back, and return it as bytes.
    
    This tool retrieves PowerPoint certificate templates from Salesforce, replaces placeholders
    with user-provided data, calculates the validity date automatically, and optionally uploads
    the modified file back to Salesforce.
    
    Workflow:
    1. Download the certificate template from Salesforce
    2. Replace <Company> with user-provided company name
    3. Replace <Tier> with user-provided tier level (Silver/Gold/Platinum)
    4. Update "Valid through" date to 31 December of the current year
    5. Upload the modified certificate back to Salesforce (if upload_back_to_salesforce is True)
    6. Return the modified file as bytes for download
    
    The tool will log upload information (ContentVersion ID and ContentDocument ID) if upload is enabled.
    
    Args:
        input_data: Contains file_id, company_name, tier, upload_back_to_salesforce, and title
        
    Returns:
        bytes: The modified certificate file as bytes for download
    """
    
    try:
        # Extract inputs
        file_id = input_data.file_id.strip()
        company_name = input_data.company_name.strip()
        tier = input_data.tier.strip()
        upload_back = input_data.upload_back_to_salesforce
        title = input_data.title.strip()
        
        # Validate file ID
        if not file_id:
            raise ValueError("file_id cannot be empty")
        
        # Get OAuth2 credentials from the connection
        creds = connections.oauth2_auth_code(SALESFORCE_APP_ID)
        base_url = creds.url.rstrip('/')
        
        # Prepare authorization headers
        headers = {
            "Authorization": f"Bearer {creds.access_token}",
            "Accept": "*/*"
        }
        
        # Track the original ContentDocument ID for creating new version
        original_content_document_id = None
        
        # Handle ContentDocument ID (069) - need to get ContentVersion ID first
        if file_id.startswith('069'):
            logger.info(f"ContentDocument ID detected: {file_id}")
            original_content_document_id = file_id
            
            # Query for the latest ContentVersion
            query = (
                f"SELECT Id FROM ContentVersion "
                f"WHERE ContentDocumentId = '{file_id}' AND IsLatest = true"
            )
            query_url = f"{base_url}/services/data/{SALESFORCE_API_VERSION}/query"
            
            query_response = requests.get(
                query_url, 
                headers=headers, 
                params={"q": query}, 
                timeout=REQUEST_TIMEOUT
            )
            query_response.raise_for_status()
            
            query_data = query_response.json()
            
            if not query_data.get('records'):
                raise ValueError(
                    f"No file found with ContentDocument ID: {file_id}"
                )
            
            # Get the ContentVersion ID
            content_version_id = query_data['records'][0]['Id']
            file_id = content_version_id
            logger.info(f"Found ContentVersion ID: {file_id}")
            
            # Set download URL for ContentVersion
            download_url = (
                f"{base_url}/services/data/{SALESFORCE_API_VERSION}/"
                f"sobjects/ContentVersion/{file_id}/VersionData"
            )
        
        # Handle ContentVersion ID (068)
        elif file_id.startswith('068'):
            logger.info(f"ContentVersion ID detected: {file_id}")
            download_url = (
                f"{base_url}/services/data/{SALESFORCE_API_VERSION}/"
                f"sobjects/ContentVersion/{file_id}/VersionData"
            )
        
        # Handle Attachment ID (00P)
        elif file_id.startswith('00P'):
            logger.info(f"Attachment ID detected: {file_id}")
            download_url = (
                f"{base_url}/services/data/{SALESFORCE_API_VERSION}/"
                f"sobjects/Attachment/{file_id}/Body"
            )
        
        else:
            raise ValueError(
                f"Unknown Salesforce ID format: {file_id}. "
                f"Expected ContentDocument (069), ContentVersion (068), "
                f"or Attachment (00P)"
            )
        
        # Download the file with streaming for large files
        logger.info(f"Downloading file from: {download_url}")
        
        response = requests.get(
            download_url,
            headers=headers,
            timeout=REQUEST_TIMEOUT,
            stream=True
        )
        
        # Raise exception for HTTP errors
        response.raise_for_status()
        
        # Get file content as bytes
        file_content = response.content
        file_size = len(file_content)
        
        logger.info(f"Successfully retrieved file ({file_size} bytes)")
        
        # Modify the PowerPoint template
        modified_content = modify_pptx_template(file_content, company_name, tier)
        
        # Generate dynamic title based on company and tier
        if title == "Partner_Plus_Certificate":
            # Create unique title with company and tier
            safe_company = company_name.replace(" ", "_").replace("/", "_")
            dynamic_title = f"Partner_Plus_Certificate_{tier}_{safe_company}"
        else:
            dynamic_title = title
        
        # Upload back to Salesforce if requested
        # NOTE: We pass None for original_content_document_id to create a NEW document
        # instead of creating a new version of the template
        if upload_back:
            upload_result = upload_file_to_salesforce(
                base_url=base_url,
                access_token=creds.access_token,
                file_content=modified_content,
                title=dynamic_title,
                original_content_document_id=None  # Always create new document
            )
            logger.info(f"Upload successful! ContentVersion ID: {upload_result.get('content_version_id')}, "
                       f"ContentDocument ID: {upload_result.get('content_document_id')}")
            
            # Log the Salesforce IDs for the agent to potentially use
            print(f"SALESFORCE_UPLOAD_INFO: ContentVersion={upload_result.get('content_version_id')}, "
                  f"ContentDocument={upload_result.get('content_document_id')}")
        
        # Return the modified file bytes for download
        return modified_content
        
    except Exception as e:
        # Return error message as bytes (text file)
        error_msg = f"Error processing file from Salesforce: {str(e)}"
        logger.error(error_msg)
        return error_msg.encode('utf-8')


# Example usage for testing
if __name__ == "__main__":
    # Test the function locally
    test_input = SalesforceFileInput(
        file_id="069NS00000VGfOjYAL",
        company_name="Tech Innovations Inc.",
        tier="Platinum",
        upload_back_to_salesforce=True,
        title="Partner_Plus_Certificate_Platinum"
    )
    
    # Note: This will fail without proper Salesforce credentials
    print("Testing Salesforce file download, modification, and upload...")
    print(f"File ID: {test_input.file_id}")
    print(f"Company: {test_input.company_name}")
    print(f"Tier: {test_input.tier}")
    print(f"Upload back: {test_input.upload_back_to_salesforce}")
    print(f"Title: {test_input.title}")